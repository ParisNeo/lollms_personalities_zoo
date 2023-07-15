from lollms.config import TypedConfig, BaseConfig, ConfigTemplate, InstallOption
from lollms.types import MSG_TYPE
from lollms.personality import APScript, AIPersonality
from lollms.paths import LollmsPaths
from lollms.helpers import ASCIIColors, trace_exception

import numpy as np
import json
from pathlib import Path
import numpy as np
import json
import subprocess

class TextVectorizer:
    def __init__(self, model, personality_config:ConfigTemplate, lollms_paths:LollmsPaths):
        
        
        self.model = model
        self.embeddings = {}
        self.texts = {}
        self.ready = False
        self.personality_config = personality_config
        self.vectorizer = None

        
        
        self.database_file = Path(lollms_paths.personal_data_path/self.personality_config["database_path"])

        self.visualize_data_at_startup=self.personality_config["visualize_data_at_startup"],
        self.visualize_data_at_add_file=self.personality_config["visualize_data_at_add_file"],
        self.visualize_data_at_generate=self.personality_config["visualize_data_at_generate"]

        if self.model.embed("hi")==None:
            self.personality_config.vectorization_method="ftidf_vectorizer"

        # Load previous state from the JSON file
        if Path(self.database_file).exists():
            ASCIIColors.success(f"Database file found : {self.database_file}")
            self.load_from_json()
            if self.visualize_data_at_startup:
                self.show_document()
            self.ready = True
        else:
            ASCIIColors.info(f"No database file found : {self.database_file}")

                
    def show_document(self, query_text=None, use_pca=True):
        import textwrap
        import seaborn as sns
        import matplotlib.pyplot as plt
        import mplcursors
        from tkinter import Tk, Text, Scrollbar, Frame, Label, TOP, BOTH, RIGHT, LEFT, Y, N, END

        
        from sklearn.manifold import TSNE
        from sklearn.decomposition import PCA
        import torch
        
        if use_pca:
            print("Showing pca representation :")
        else:
            print("Showing t-sne representation :")
        texts = list(self.texts.values())
        embeddings = self.embeddings
        if len(embeddings.values())>=2:
            # Normalize embeddings
            norms = np.linalg.norm(embeddings, axis=1)
            normalized_embeddings = embeddings / norms[:, np.newaxis]

            # Embed the query text
            query_embedding = self.embed_query(query_text)
            query_embedding = query_embedding.detach().squeeze().numpy()
            query_normalized_embedding = query_embedding / np.linalg.norm(query_embedding)

            # Combine the query embedding with the document embeddings
            combined_embeddings = np.vstack((normalized_embeddings, query_normalized_embedding))

            if use_pca:
                # Use PCA for dimensionality reduction
                pca = PCA(n_components=2)
                embeddings_2d = pca.fit_transform(combined_embeddings)
            else:
                # Use t-SNE for dimensionality reduction
                # Adjust the perplexity value
                perplexity = min(30, combined_embeddings.shape[0] - 1)
                tsne = TSNE(n_components=2, perplexity=perplexity)
                embeddings_2d = tsne.fit_transform(combined_embeddings)


            # Create a scatter plot using Seaborn
            sns.scatterplot(x=embeddings_2d[:-1, 0], y=embeddings_2d[:-1, 1])  # Plot document embeddings
            plt.scatter(embeddings_2d[-1, 0], embeddings_2d[-1, 1], color='red')  # Plot query embedding

            # Add labels to the scatter plot
            for i, (x, y) in enumerate(embeddings_2d[:-1]):
                plt.text(x, y, str(i), fontsize=8)

            plt.xlabel('Dimension 1')
            plt.ylabel('Dimension 2')
            if use_pca:      
                plt.title('Embeddings Scatter Plot based on PCA')
            else:
                plt.title('Embeddings Scatter Plot based on t-SNE')
            # Enable mplcursors to show tooltips on hover
            cursor = mplcursors.cursor(hover=True)

            # Define the hover event handler
            @cursor.connect("add")
            def on_hover(sel):
                index = sel.target.index
                if index > 0:
                    text = texts[index]
                    wrapped_text = textwrap.fill(text, width=50)  # Wrap the text into multiple lines
                    sel.annotation.set_text(f"Index: {index}\nText:\n{wrapped_text}")
                else:
                    sel.annotation.set_text("Query")

            # Define the click event handler using matplotlib event handling mechanism
            def on_click(event):
                if event.xdata is not None and event.ydata is not None:
                    x, y = event.xdata, event.ydata
                    distances = ((embeddings_2d[:, 0] - x) ** 2 + (embeddings_2d[:, 1] - y) ** 2)
                    index = distances.argmin()
                    text = texts[index] if index < len(texts) else query_text

                    # Open a new Tkinter window with the content of the text
                    root = Tk()
                    root.title(f"Text for Index {index}")
                    frame = Frame(root)
                    frame.pack(fill=BOTH, expand=True)

                    label = Label(frame, text="Text:")
                    label.pack(side=TOP, padx=5, pady=5)

                    text_box = Text(frame)
                    text_box.pack(side=TOP, padx=5, pady=5, fill=BOTH, expand=True)
                    text_box.insert(END, text)

                    scrollbar = Scrollbar(frame)
                    scrollbar.pack(side=RIGHT, fill=Y)
                    scrollbar.config(command=text_box.yview)
                    text_box.config(yscrollcommand=scrollbar.set)

                    text_box.config(state="disabled")

                    root.mainloop()

            # Connect the click event handler to the figure
            plt.gcf().canvas.mpl_connect("button_press_event", on_click)
            plt.show()
        
    def index_document(self, document_id, text, chunk_size, overlap_size, force_vectorize=False):

        if document_id in self.embeddings and not force_vectorize:
            print(f"Document {document_id} already exists. Skipping vectorization.")
            return

        # Tokenize text
        tokens = self.model.tokenize(text)

        # Split tokens into sentences
        sentences = self.model.detokenize(tokens).split('. ')

        # Generate chunks with overlap and sentence boundaries
        chunks = []
        current_chunk = []
        for sentence in sentences:
            sentence_tokens = self.model.tokenize(sentence)
            if len(current_chunk) + len(sentence_tokens) <= chunk_size:
                current_chunk.extend(sentence_tokens)
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence_tokens

        if current_chunk:
            chunks.append(current_chunk)

        # Generate overlapping chunks
        overlapping_chunks = []
        for i in range(len(chunks)):
            chunk_start = i * (chunk_size - overlap_size)
            chunk_end = min(chunk_start + chunk_size, len(tokens))
            chunk = tokens[chunk_start:chunk_end]
            overlapping_chunks.append(chunk)


        if self.personality_config.vectorization_method=="ftidf_vectorizer":
            from sklearn.feature_extraction.text import TfidfVectorizer
            self.vectorizer = TfidfVectorizer()
            data = [self.model.detokenize(chunk) for chunk in overlapping_chunks]
            self.vectorizer.fit(data)

        # Generate embeddings for each chunk
        for i, chunk in enumerate(overlapping_chunks):
            # Store chunk ID, embedding, and original text
            chunk_id = f"{document_id}_chunk_{i + 1}"
            self.texts[chunk_id] = self.model.detokenize(chunk[:chunk_size])
            if self.personality_config.vectorization_method=="ftidf_vectorizer":
                self.embeddings[chunk_id] = self.vectorizer.transform([self.texts[chunk_id]]).toarray()
            else:
                self.embeddings[chunk_id] = self.model.embed(self.texts[chunk_id])

        self.save_to_json()
        self.ready = True
        if self.visualize_data_at_add_file:
            self.show_document()


    def embed_query(self, query_text):
        # Generate query embedding
        if self.personality_config.vectorization_method=="ftidf_vectorizer":
            query_embedding = self.vectorizer.transform(query_text).toarray()
        else:
            query_embedding = self.model.embed(query_text)

        return query_embedding

    def recover_text(self, query_embedding, top_k=1):
        from sklearn.metrics.pairwise import cosine_similarity
        similarities = {}
        for chunk_id, chunk_embedding in self.embeddings.items():
            similarity = cosine_similarity(query_embedding, chunk_embedding)[0][0]
            similarities[chunk_id] = similarity

        # Sort the similarities and retrieve the top-k most similar embeddings
        sorted_similarities = sorted(similarities.items(), key=lambda x: x[1], reverse=True)[:top_k]

        # Retrieve the original text associated with the most similar embeddings
        texts = [self.texts[chunk_id] for chunk_id, _ in sorted_similarities]

        if self.visualize_data_at_generate:
            self.show_document()

        return texts

    def save_to_json(self):
        state = {
            "embeddings": {str(k): v.tolist() for k, v in self.embeddings.items()},
            "texts": self.texts,
        }
        with open(self.database_file, "w") as f:
            json.dump(state, f)

    def load_from_json(self):

        ASCIIColors.info("Loading vectorized documents")
        with open(self.database_file, "r") as f:
            state = json.load(f)
            self.embeddings = {k: v for k, v in state["embeddings"].items()}
            self.texts = state["texts"]
            self.ready = True
        if self.vector_store and self.personality_config.vectorization_method=="ftidf_vectorizer":
            from sklearn.feature_extraction.text import TfidfVectorizer
            data = list(self.texts.values())
            if len(data)>0:
                self.vectorizer = TfidfVectorizer()
                self.vectorizer.fit(data)

class Processor(APScript):
    """
    A class that processes model inputs and outputs.

    Inherits from APScript.
    """

    def __init__(
                 self, 
                 personality: AIPersonality
                ) -> None:
        
        self.word_callback = None    

        personality_config_template = ConfigTemplate(
            [
                {"name":"vectorization_method","type":"str","value":f"model_embedding", "options":["model_embedding", "ftidf_vectorizer"], "help":"Vectoriazation method to be used (changing this should reset database)"},
                {"name":"database_path","type":"str","value":f"{personality.name}_db.json", "help":"Path to the database"},
                {"name":"max_chunk_size","type":"int","value":512, "min":10, "max":personality.config["ctx_size"],"help":"Maximum size of text chunks to vectorize"},
                {"name":"chunk_overlap","type":"int","value":20, "min":0, "max":personality.config["ctx_size"],"help":"Overlap between chunks"},
                
                {"name":"max_answer_size","type":"int","value":512, "min":10, "max":personality.config["ctx_size"],"help":"Maximum number of tokens to allow the generator to generate as an answer to your question"},
                
                {"name":"visualize_data_at_startup","type":"bool","value":False, "help":"If true, the database will be visualized at startup"},
                {"name":"visualize_data_at_add_file","type":"bool","value":False, "help":"If true, the database will be visualized when a new file is added"},
                {"name":"visualize_data_at_generate","type":"bool","value":False, "help":"If true, the database will be visualized at generation time"},
            ]
            )
        personality_config_vals = BaseConfig.from_template(personality_config_template)

        personality_config = TypedConfig(
            personality_config_template,
            personality_config_vals
        )
        super().__init__(
                            personality,
                            personality_config,
                            [
                                {
                                    "name": "idle",
                                    "commands": { # list of commands
                                        "help":self.help,
                                        "set_database": self.set_database,
                                        "clear_database": self.clear_database
                                    },
                                    "default": self.chat_with_doc
                                },                           
                            ]
                        )
        self.state = 0
        self.ready = False
        self.personality = personality
        self.callback = None
        self.vector_store = None


    def install(self):
        super().install()
        # Get the current directory
        root_dir = self.personality.lollms_paths.personal_path
        # We put this in the shared folder in order as this can be used by other personalities.
        shared_folder = root_dir/"shared"
        sd_folder = shared_folder / "sd"

        requirements_file = self.personality.personality_package_path / "requirements.txt"
        # Step 2: Install dependencies using pip from requirements.txt
        subprocess.run(["pip", "install", "--upgrade", "-r", str(requirements_file)])            
        try:
            print("Checking pytorch")
            import torch
            import torchvision
            if torch.cuda.is_available():
                print("CUDA is supported.")
            else:
                print("CUDA is not supported. Reinstalling PyTorch with CUDA support.")
                self.reinstall_pytorch_with_cuda()
        except Exception as ex:
            self.reinstall_pytorch_with_cuda()

        # Step 1: Clone repository
        if not sd_folder.exists():
            subprocess.run(["git", "clone", "https://github.com/CompVis/stable-diffusion.git", str(sd_folder)])

        # Step 2: Install the Python package inside sd folder
        subprocess.run(["pip", "install", "--upgrade", str(sd_folder)])

        # Step 3: Create models/Stable-diffusion folder if it doesn't exist
        models_folder = shared_folder / "sd_models"
        models_folder.mkdir(parents=True, exist_ok=True)

        ASCIIColors.success("Installed successfully")

    

    def help(self, prompt):
        self.full(self.personality.help, self.callback)

    def set_database(self, prompt):
        self.goto_state("waiting_for_file")

    def clear_database(self,prompt):
        self.vector_store.clear_database()

    def chat_with_doc(self, prompt):
        self.step_start("Recovering data")
        ASCIIColors.blue("Recovering data")
        if self.vector_store.ready:
            docs = self.vector_store.recover_text(self.vector_store.embed_query(prompt), top_k=3)
            docs = '\n'.join([f"Doc{i}:\n{v}" for i,v in enumerate(docs)])
            full_text = self.personality.personality_conditioning+"\n### Docs:\n"+docs+"\n### Question: "+prompt+"\n### Answer:"
            ASCIIColors.blue("-------------- Documentation -----------------------")
            ASCIIColors.blue(full_text)
            ASCIIColors.blue("----------------------------------------------------")
            ASCIIColors.blue("Thinking")
            self.step_end("Recovering data")
            self.step_start("Thinking",self.callback)
            output = self.generate(full_text, self.personality_config["max_answer_size"])
            self.step_end("Thinking",self.callback)
            self.full(output, self.callback)
        else:
            self.full("Vector store is not ready. Please send me a document to use. Use Send file command form your chatbox menu to trigger this.", self.callback)

    @staticmethod        
    def read_pdf_file(file_path):
        import PyPDF2
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text

    @staticmethod
    def read_docx_file(file_path):
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        import csv
        with open(file_path, 'r') as file:
            csv_reader = csv.reader(file)
            lines = [row for row in csv_reader]
        return lines    

    @staticmethod
    def read_html_file(file_path):
        from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    @staticmethod
    def read_pptx_file(file_path):
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content
    
    def build_db(self):
        if self.vector_store is None:
            self.vector_store = TextVectorizer(
                                        self.personality.model,
                                        self.personality_config,
                                        self.personality.lollms_paths
                                    )        
        if len(self.vector_store.embeddings)>0:
            self.ready = True

        ASCIIColors.info("-> Vectorizing the database"+ASCIIColors.color_orange)
        if self.callback is not None:
            self.callback("Vectorizing the database", MSG_TYPE.MSG_TYPE_CHUNK)
        for file in self.files:
            try:
                if Path(file).suffix==".pdf":
                    text =  Processor.read_pdf_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_docx_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_pptx_file(file)
                elif Path(file).suffix==".json":
                    text =  Processor.read_json_file(file)
                elif Path(file).suffix==".csv":
                    text =  Processor.read_csv_file(file)
                elif Path(file).suffix==".html":
                    text =  Processor.read_html_file(file)
                else:
                    text =  Processor.read_text_file(file)
                try:
                    chunk_size=int(self.personality_config["max_chunk_size"])
                except:
                    ASCIIColors.warning(f"Couldn't read chunk size. Verify your configuration file")
                    chunk_size=512
                try:
                    overlap_size=int(self.personality_config["chunk_overlap"])
                except:
                    ASCIIColors.warning(f"Couldn't read chunk size. Verify your configuration file")
                    overlap_size=50

                self.vector_store.index_document(file, text, chunk_size=chunk_size, overlap_size=overlap_size)
                
                print(ASCIIColors.color_reset)
                ASCIIColors.success(f"File {file} vectorized successfully")
                self.ready = True
            except Exception as ex:
                ASCIIColors.error(f"Couldn't vectorize {file}: The vectorizer threw this exception:{ex}")
                trace_exception(ex)

    def add_file(self, path):
        super().add_file(path)
        self.prepare()
        try:
            self.step_start("Vectorizing database",self.callback)
            self.build_db()
            self.step_end("Vectorizing database",self.callback)
            self.ready = True
            return True
        except Exception as ex:
            ASCIIColors.error(f"Couldn't vectorize the database: The vectgorizer threw this exception: {ex}")
            trace_exception(ex)
            return False        

    def prepare(self):
        if self.vector_store is None:
            self.vector_store = TextVectorizer(
                                        self.personality.model,
                                        self.personality_config,
                                        self.personality.lollms_paths
                                    )    

        if self.vector_store and self.personality_config.vectorization_method=="ftidf_vectorizer":
            from sklearn.feature_extraction.text import TfidfVectorizer
            data = list(self.vector_store.texts.values())
            if len(data)>0:
                self.vectorizer = TfidfVectorizer()
                self.vectorizer.fit(data)

        if len(self.vector_store.embeddings)>0:
            self.ready = True

    def run_workflow(self, prompt, previous_discussion_text="", callback=None):
        """
        Runs the workflow for processing the model input and output.

        This method should be called to execute the processing workflow.

        Args:
            generate_fn (function): A function that generates model output based on the input prompt.
                The function should take a single argument (prompt) and return the generated text.
            prompt (str): The input prompt for the model.
            previous_discussion_text (str, optional): The text of the previous discussion. Default is an empty string.

        Returns:
            None
        """
        # State machine
        self.callback = callback
        self.prepare()

        self.process_state(prompt, callback)

        return ""



